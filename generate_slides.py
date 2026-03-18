#!/usr/bin/env python3
"""Generate BCM IaC PPTX presentation from slide content."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD   = RGBColor(0x22, 0x22, 0x3A)
GREEN     = RGBColor(0x76, 0xB9, 0x00)
CYAN      = RGBColor(0x00, 0xBC, 0xD4)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY     = RGBColor(0xCC, 0xCC, 0xCC)
ORANGE    = RGBColor(0xFF, 0x98, 0x00)
RED       = RGBColor(0xEF, 0x53, 0x50)

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)
    tx = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(44); p.font.color.rgb = WHITE; p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph(); p2.text = subtitle
        p2.font.size = Pt(20); p2.font.color.rgb = LGRAY
        p2.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(title, bullets, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    # Title bar
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = tx.text_frame; p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(32); p.font.color.rgb = GREEN; p.font.bold = True
    # Green underline
    ln = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(12), Pt(3))  # rectangle
    ln.fill.solid(); ln.fill.fore_color.rgb = GREEN; ln.line.fill.background()
    # Bullets
    tx2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.5), Inches(5.5))
    tf2 = tx2.text_frame; tf2.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        if b.startswith("##"):
            p.text = b.lstrip("# "); p.font.size = Pt(20); p.font.color.rgb = CYAN; p.font.bold = True
            p.space_before = Pt(12)
        elif b.startswith("```"):
            p.text = b.strip("`"); p.font.size = Pt(14); p.font.color.rgb = ORANGE
            p.font.name = "Consolas"
        elif b.startswith("|"):
            p.text = b; p.font.size = Pt(13); p.font.color.rgb = LGRAY
            p.font.name = "Consolas"
        else:
            p.text = b; p.font.size = Pt(16); p.font.color.rgb = WHITE
        p.space_after = Pt(4)
    if note:
        p = tf2.add_paragraph(); p.text = note
        p.font.size = Pt(13); p.font.color.rgb = LGRAY; p.font.italic = True
    return slide

def add_table_slide(title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = tx.text_frame; p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(32); p.font.color.rgb = GREEN; p.font.bold = True
    
    cols = len(headers)
    tbl = slide.shapes.add_table(len(rows)+1, cols, Inches(0.5), Inches(1.3), Inches(12), Inches(5.5)).table
    
    # Adjust column widths
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i); cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.color.rgb = WHITE; p.font.bold = True
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x30,0x30,0x50)
    
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r+1, c); cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12); p.font.color.rgb = LGRAY
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD if r % 2 == 0 else RGBColor(0x28,0x28,0x42)
    return slide

# ─── SLIDE 1: Title ─────────────────────────────────────────────────
add_title_slide(
    "BCM Infrastructure as Code",
    "Multi-Environment Ansible Automation for BCM 11.0\n\nOne Repo • 6 Environments • 11 Playbooks • Explicit Commands"
)

# ─── SLIDE 2: Environments ──────────────────────────────────────────
add_table_slide("6 Environments — One Repo", 
    ["Environment", "What It Is", "Scale"],
    [
        ["local-lab", "Your workstation — KVM VMs", "1 head + 3 nodes"],
        ["az2-staging", "AZ2 staging cluster", "1 CPU head + 1 GPU + few CPU"],
        ["az2-staging-bmaas", "BMaaS staging cluster", "1 head + 3 CPU + 1 GPU"],
        ["az2-bmaas-prod", "BMaaS production", "1 head + 66 DGX B200"],
        ["az1-prod", "AZ1 production", "TBD"],
        ["az2-prod", "AZ2 production", "TBD"],
    ]
)

# ─── SLIDE 3: Repo Structure ────────────────────────────────────────
add_content_slide("Repository Structure", [
    "## WHERE — inventories/",
    "  Per-environment: hosts.yml + group_vars/all.yml",
    "  local-lab/ • az2-staging/ • az2-staging-bmaas/ • az2-bmaas-prod/ • az1-prod/ • az2-prod/",
    "",
    "## WHAT — playbooks/  (11 playbooks)",
    "  cluster-health • dns-config • nodegroup-config • firmware-audit • firmware-upgrade",
    "  node-lifecycle • slurm-ops • slurm-jobs • debug-bundle • bcm-config • bcm-image-mgmt",
    "",
    "## HOW — roles/  (4 reusable roles)",
    "  bcm-dns • bcm-nodegroup • bcm-health • bcm-slurm",
    "",
    "## DOCS — docs/  (7 guides)",
    "  Runbook • Concepts • DNS best practices • Operations • Install troubleshooting",
])

# ─── SLIDE 4: How to Run ────────────────────────────────────────────
add_content_slide("How to Run — Explicit Commands", [
    "## Pattern:",
    "ansible-playbook -i inventories/<ENV>/hosts.yml playbooks/<PLAYBOOK>.yml",
    "",
    "## Examples:",
    "ansible-playbook -i inventories/local-lab/hosts.yml playbooks/cluster-health.yml",
    "ansible-playbook -i inventories/az2-bmaas-prod/hosts.yml playbooks/dns-config.yml",
    "ansible-playbook -i inventories/az2-staging/hosts.yml playbooks/firmware-audit.yml",
    "",
    "## Dry-run (add --check --diff):",
    "ansible-playbook -i inventories/az2-bmaas-prod/hosts.yml playbooks/dns-config.yml --check --diff",
    "",
    "## Specific tags only:",
    "ansible-playbook -i inventories/az2-bmaas-prod/hosts.yml playbooks/firmware-audit.yml --tags gpu,nic",
])

# ─── SLIDE 5: All 11 Playbooks ──────────────────────────────────────
add_table_slide("All 11 Playbooks", 
    ["Playbook", "What It Does"],
    [
        ["cluster-health.yml", "Check services, disk, memory, Slurm status"],
        ["dns-config.yml", "Set nameservers, search domains via cmsh"],
        ["nodegroup-config.yml", "Manage categories, overlays, node assignments"],
        ["firmware-audit.yml", "Collect BIOS, BMC, GPU, NIC firmware versions"],
        ["firmware-upgrade.yml", "Upgrade firmware (dry-run default, safety guards)"],
        ["node-lifecycle.yml", "Power, provision, drain, resume nodes"],
        ["slurm-ops.yml", "Slurm partition management"],
        ["slurm-jobs.yml", "Submit and monitor Slurm jobs"],
        ["debug-bundle.yml", "Collect diagnostic logs and GPU info"],
        ["bcm-config.yml", "Audit BCM configuration objects"],
        ["bcm-image-mgmt.yml", "Audit and manage software images"],
    ]
)

# ─── SLIDE 6: Day-0 Operations ──────────────────────────────────────
add_table_slide("Day-0 Operations — Initial Setup",
    ["#", "Operation", "Playbook", "Status"],
    [
        ["1", "Health baseline", "cluster-health.yml", "✅ Tested"],
        ["2", "BCM config audit", "bcm-config.yml", "✅ Tested"],
        ["3", "Image audit", "bcm-image-mgmt.yml", "✅ Tested"],
        ["4", "Configure DNS", "dns-config.yml", "🔲 Pending"],
        ["5", "Setup node categories", "nodegroup-config.yml", "🔲 Pending"],
        ["6", "Provision nodes", "node-lifecycle.yml", "✅ Tested"],
        ["7", "Setup Slurm", "slurm-ops.yml", "✅ Tested"],
        ["8", "Firmware audit", "firmware-audit.yml", "🔲 Pending"],
        ["9", "Firmware upgrade", "firmware-upgrade.yml", "🔲 Pending"],
    ]
)

# ─── SLIDE 7: Day-2 Operations ──────────────────────────────────────
add_table_slide("Day-2 Operations — Ongoing Management",
    ["#", "Operation", "Playbook", "Frequency"],
    [
        ["1", "Health check", "cluster-health.yml", "Daily"],
        ["2", "Debug bundle", "debug-bundle.yml", "On incident"],
        ["3", "Node power cycle", "node-lifecycle.yml", "As needed"],
        ["4", "Node reprovision", "node-lifecycle.yml", "After RMA"],
        ["5", "Slurm job monitoring", "slurm-jobs.yml", "Daily"],
        ["6", "Image update", "bcm-image-mgmt.yml", "On patch"],
        ["7", "DNS change", "dns-config.yml", "On network change"],
        ["8", "Category change", "nodegroup-config.yml", "On scale"],
        ["9", "Firmware upgrade", "firmware-upgrade.yml", "Quarterly"],
        ["10", "Firmware audit", "firmware-audit.yml", "Monthly"],
    ]
)

# ─── SLIDE 8: Firmware Audit ────────────────────────────────────────
add_content_slide("Firmware Audit — firmware-audit.yml", [
    "## Collects from every node:",
    "  • BIOS version         → dmidecode -t bios",
    "  • BMC firmware          → ipmitool mc info",
    "  • GPU driver + VBIOS    → nvidia-smi --query-gpu",
    "  • NIC firmware          → mlxfwmanager --query",
    "  • BCM + Slurm versions  → cmsh get version",
    "",
    "## Run:",
    "ansible-playbook -i inventories/az2-bmaas-prod/hosts.yml playbooks/firmware-audit.yml",
    "",
    "## GPU only:",
    "ansible-playbook -i inventories/az2-bmaas-prod/hosts.yml playbooks/firmware-audit.yml --tags gpu",
    "",
    "## Report saved to: /tmp/firmware-audit-<env>/firmware-audit-<timestamp>.md",
])

# ─── SLIDE 9: Firmware Upgrade ───────────────────────────────────────
add_content_slide("Firmware Upgrade — Safety-First", [
    "## firmware-upgrade.yml — DRY-RUN by default",
    "",
    "  1. Pre-flight    → checks Slurm, aborts if jobs running",
    "  2. Drain         → removes nodes from Slurm",
    "  3. Upgrade GPU   → updates driver in BCM image",
    "  4. Upgrade BIOS  → flashes BIOS firmware",
    "  5. Upgrade BMC   → updates BMC firmware",
    "  6. Upgrade NIC   → flashes Mellanox ConnectX",
    "  7. Reboot        → reboots nodes, waits 120s",
    "  8. Verify        → re-checks all versions",
    "  9. Resume        → puts nodes back in Slurm",
    "",
    "## Apply GPU driver upgrade:",
    "ansible-playbook -i inventories/az2-bmaas-prod/hosts.yml playbooks/firmware-upgrade.yml -e fw_apply=true -e target_gpu_driver=560",
])

# ─── SLIDE 10: Image vs Overlay ─────────────────────────────────────
add_table_slide("Image vs Overlay — What Goes Where",
    ["✅ In the Image", "✅ In the Overlay"],
    [
        ["OS base packages", "/etc/resolv.conf (DNS)"],
        ["NVIDIA GPU drivers", "NTP servers"],
        ["CUDA toolkit", "LDAP client config"],
        ["Container runtime", "Mount points (/scratch)"],
        ["Slurm client (slurmd)", "SSH authorized_keys"],
        ["Monitoring agents", "Kernel parameters"],
        ["System libs", "Network bonding"],
    ]
)

# ─── SLIDE 11: Promotion Flow ───────────────────────────────────────
add_content_slide("Promotion Flow", [
    "## local-lab → az2-staging → az2-staging-bmaas → az2-bmaas-prod",
    "",
    "  1. DEVELOP on local-lab         (KVM VMs on your workstation)",
    "  2. TEST on az2-staging           (real hardware, small scale)",
    "  3. VALIDATE on az2-staging-bmaas (BMaaS-specific testing)",
    "  4. DEPLOY to az2-bmaas-prod      (production, 66 B200 nodes)",
    "",
    "## All on main branch — no branch-per-environment",
    "",
    "## Same playbooks, different inventories:",
    "ansible-playbook -i inventories/local-lab/hosts.yml playbooks/dns-config.yml",
    "ansible-playbook -i inventories/az2-staging/hosts.yml playbooks/dns-config.yml",
    "ansible-playbook -i inventories/az2-bmaas-prod/hosts.yml playbooks/dns-config.yml",
], note="Same code → different variables → different outcomes")

# ─── SLIDE 12: Where to Find Things ─────────────────────────────────
add_table_slide("Where to Find Things",
    ["I want to...", "Look at..."],
    [
        ["See environment config", "inventories/<env>/group_vars/all.yml"],
        ["See what a playbook does", "playbooks/<name>.yml"],
        ["See reusable logic", "roles/<name>/tasks/main.yml"],
        ["Dry-run something", "Append --check --diff"],
        ["Run specific tags only", "Append --tags <tag1>,<tag2>"],
        ["BCM concepts", "docs/bcm-concepts-guide.md"],
        ["DNS best practices", "docs/bcm-dns-and-image-best-practices.md"],
        ["Install troubleshooting", "docs/bcm-install-troubleshooting.md"],
        ["Copy-paste commands", "docs/tested-runbook.md"],
    ]
)

# ─── SLIDE 13: Testing Backlog ───────────────────────────────────────
add_table_slide("Testing Backlog",
    ["Priority", "Playbook", "Status"],
    [
        ["P0", "cluster-health.yml", "✅ Tested on local-lab"],
        ["P0", "dns-config.yml", "🔲 Need nodes UP"],
        ["P0", "nodegroup-config.yml", "🔲 Need nodes UP"],
        ["P1", "node-lifecycle.yml", "✅ Tested on local-lab"],
        ["P1", "slurm-ops.yml", "✅ Tested on local-lab"],
        ["P1", "debug-bundle.yml", "✅ Tested on local-lab"],
        ["P1", "firmware-audit.yml", "🔲 Need nodes UP"],
        ["P1", "firmware-upgrade.yml", "🔲 Need nodes UP"],
        ["P2", "bcm-image-mgmt.yml", "✅ Tested on local-lab"],
        ["P2", "bcm-config.yml", "✅ Tested on local-lab"],
        ["P2", "slurm-jobs.yml", "✅ Tested on local-lab"],
    ]
)

# ─── SLIDE 14: What's Next ──────────────────────────────────────────
add_content_slide("What's Next", [
    "## Planned Playbooks:",
    "  • burn-in-test.yml    — post-RMA node validation",
    "  • ldap-config.yml     — LDAP/authentication setup",
    "  • network-bonding.yml — IB/ETH bond configuration",
    "",
    "## Infrastructure:",
    "  • CI/CD pipeline (lint on PR + dry-run on merge)",
    "  • Ansible Vault for secrets management",
    "  • GPU-specific playbooks (DCGM, nvidia-smi health)",
    "",
    "## Immediate:",
    "  • Get local-lab nodes fully provisioned (PXE booting now)",
    "  • Test DNS + nodegroup playbooks on live nodes",
    "  • Fill in az1-prod and az2-prod with real IPs/hostnames",
])

outpath = os.path.join(os.path.dirname(os.path.abspath(__file__)), "BCM-IaC-Walkthrough.pptx")
prs.save(outpath)
print(f"Saved: {outpath}")
print(f"Slides: {len(prs.slides)}")
